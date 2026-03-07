"""
JARVIS API Server — FastAPI application exposing the AI Second Brain.

Endpoints:
    Auth:       POST /auth/register, POST /auth/login
    Documents:  POST /documents/upload, GET /documents, GET /documents/{id}, DELETE /documents/{id}
    Knowledge:  POST /knowledge/query
    Outputs:    GET /flashcards, GET /quizzes, GET /study-plans
    System:     GET /health, GET /status
"""

from __future__ import annotations

import hashlib
import io
import os
import time
import uuid
from contextlib import asynccontextmanager
from datetime import datetime, timedelta, timezone
from typing import Any, Dict, List, Optional

from fastapi import (
    Depends,
    FastAPI,
    File,
    HTTPException,
    Request,
    UploadFile,
    status,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel, Field

from app.graph.workflow import JarvisWorkflow, build_workflow
from app.utils.config import settings
from app.utils.logger import get_logger

logger = get_logger("api")


# ═══════════════════════════════════════════════════════════════════════════════
# JWT helpers (lightweight — avoids pulling in python-jose at import time)
# ═══════════════════════════════════════════════════════════════════════════════

try:
    from jose import JWTError, jwt as jose_jwt
    _HAS_JOSE = True
except ImportError:
    _HAS_JOSE = False


def _create_access_token(data: dict, expires_delta: Optional[timedelta] = None) -> str:
    if not _HAS_JOSE:
        # Fallback: create a simple opaque token (dev-only)
        return str(uuid.uuid4())
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + (expires_delta or timedelta(minutes=settings.app.jwt_expiry_minutes))
    to_encode.update({"exp": expire})
    return jose_jwt.encode(to_encode, settings.app.secret_key, algorithm=settings.app.jwt_algorithm)


def _decode_token(token: str) -> dict:
    if not _HAS_JOSE:
        raise HTTPException(status_code=401, detail="JWT library not installed; auth disabled")
    try:
        return jose_jwt.decode(token, settings.app.secret_key, algorithms=[settings.app.jwt_algorithm])
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid or expired token")


# ═══════════════════════════════════════════════════════════════════════════════
# Application lifespan (init / shutdown)
# ═══════════════════════════════════════════════════════════════════════════════

workflow: Optional[JarvisWorkflow] = None


@asynccontextmanager
async def lifespan(app: FastAPI):
    global workflow
    logger.info("Starting JARVIS API …", event_type="api_startup")
    workflow = build_workflow()
    await workflow.initialize()
    logger.info("JARVIS API ready", event_type="api_ready")
    yield
    logger.info("Shutting down JARVIS API …", event_type="api_shutdown")
    if workflow:
        await workflow.shutdown()


# ═══════════════════════════════════════════════════════════════════════════════
# FastAPI app
# ═══════════════════════════════════════════════════════════════════════════════

app = FastAPI(
    title="JARVIS AI Second Brain",
    version="1.0.0",
    description="Multi-agent AI personal knowledge manager powered by Azure OpenAI",
    lifespan=lifespan,
)

# CORS — allow all origins in development
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ═══════════════════════════════════════════════════════════════════════════════
# Dependency: extract current user from Authorization header
# ═══════════════════════════════════════════════════════════════════════════════

async def get_current_user(request: Request) -> dict:
    """
    Extract and validate the JWT bearer token.

    Falls back to a default dev user when JWT is not configured.
    """
    auth = request.headers.get("Authorization", "")
    if not auth.startswith("Bearer "):
        if settings.app.debug:
            return {"user_id": "dev_user", "email": "dev@jarvis.local"}
        raise HTTPException(status_code=401, detail="Missing Authorization header")

    token = auth.split(" ", 1)[1]
    payload = _decode_token(token)
    return {"user_id": payload.get("sub", ""), "email": payload.get("email", "")}


# ═══════════════════════════════════════════════════════════════════════════════
# Pydantic request / response models
# ═══════════════════════════════════════════════════════════════════════════════

class RegisterRequest(BaseModel):
    name: str = Field(..., min_length=1, max_length=100)
    email: str = Field(..., min_length=5, max_length=200)
    password: str = Field(..., min_length=6, max_length=128)


class LoginRequest(BaseModel):
    email: str
    password: str


class QueryRequest(BaseModel):
    query: str = Field(..., min_length=1, max_length=4000)
    session_id: str = Field(default="")


class QueryResponse(BaseModel):
    status: str
    request_id: str
    response: dict
    metadata: dict


class HealthResponse(BaseModel):
    status: str
    version: str
    uptime_seconds: float
    services: dict


_start_time = time.time()


# ═══════════════════════════════════════════════════════════════════════════════
# Routes — Health & Status
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/health", response_model=HealthResponse, tags=["System"])
async def health_check():
    """System health and service status."""
    return HealthResponse(
        status="healthy",
        version="1.0.0",
        uptime_seconds=round(time.time() - _start_time, 2),
        services=settings.validate_azure_services(),
    )


@app.get("/status", tags=["System"])
async def system_status():
    """Extended system status including configuration."""
    return {
        "app_name": settings.app.app_name,
        "debug": settings.app.debug,
        "database": settings.database.database_url.split("///")[0] + "///***",
        "services": settings.validate_azure_services(),
        "workflow_ready": workflow is not None and workflow._initialized,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# Routes — Authentication
# ═══════════════════════════════════════════════════════════════════════════════

def _hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()


@app.post("/auth/register", tags=["Auth"], status_code=201)
async def register(body: RegisterRequest):
    """Register a new user account."""
    existing = await workflow.memory.get_user_by_email(body.email)
    if existing:
        raise HTTPException(status_code=409, detail="Email already registered")

    pw_hash = _hash_password(body.password)
    user = await workflow.memory.create_user(body.name, body.email, pw_hash)

    token = _create_access_token({"sub": user["id"], "email": body.email})

    return {
        "message": "User registered successfully",
        "user_id": user["id"],
        "token": token,
    }


@app.post("/auth/login", tags=["Auth"])
async def login(body: LoginRequest):
    """Authenticate and receive a JWT token."""
    user = await workflow.memory.get_user_by_email(body.email)
    if not user:
        raise HTTPException(status_code=401, detail="Invalid credentials")

    if user.get("password_hash") != _hash_password(body.password):
        raise HTTPException(status_code=401, detail="Invalid credentials")

    token = _create_access_token({"sub": user["id"], "email": body.email})

    return {
        "message": "Login successful",
        "user_id": user["id"],
        "token": token,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# Routes — Knowledge Query (core endpoint)
# ═══════════════════════════════════════════════════════════════════════════════

@app.post("/knowledge/query", response_model=QueryResponse, tags=["Knowledge"])
async def query_knowledge(body: QueryRequest, user: dict = Depends(get_current_user)):
    """
    Process a natural-language query through the full agent pipeline.

    This is the primary endpoint — it invokes the LangGraph workflow
    (safety → retrieve → plan → decompose → action_plan → execute → learn).
    """
    result = await workflow.run(
        user_input=body.query,
        user_id=user["user_id"],
        session_id=body.session_id,
    )

    return QueryResponse(
        status=result["status"],
        request_id=result["request_id"],
        response=result["response"],
        metadata=result["metadata"],
    )


# ═══════════════════════════════════════════════════════════════════════════════
# Routes — Document Management
# ═══════════════════════════════════════════════════════════════════════════════

@app.post("/documents/upload", tags=["Documents"], status_code=201)
async def upload_document(
    file: UploadFile = File(...),
    user: dict = Depends(get_current_user),
):
    """
    Upload a document for ingestion into the knowledge base.

    Supported types: .pdf, .docx, .txt, .md, .pptx
    The document is chunked, embedded, and stored in the vector index.
    """
    allowed_extensions = {".pdf", ".docx", ".txt", ".md", ".pptx"}
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in allowed_extensions:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{ext}'. Allowed: {', '.join(allowed_extensions)}",
        )

    max_bytes = settings.app.max_file_size_mb * 1024 * 1024
    content = await file.read()
    if len(content) > max_bytes:
        raise HTTPException(
            status_code=413,
            detail=f"File exceeds {settings.app.max_file_size_mb}MB limit",
        )

    # Register document in DB
    doc = await workflow.memory.create_document(
        user_id=user["user_id"],
        filename=file.filename or "untitled",
        file_type=ext,
    )

    # Simple text extraction (production would use Azure Document Intelligence)
    try:
        text = _extract_text(content, ext)
        chunks = _chunk_text(text, settings.app.chunk_size, settings.app.chunk_overlap)

        await workflow.memory.store_document_chunks(
            document_id=doc["id"],
            user_id=user["user_id"],
            filename=file.filename or "untitled",
            chunks=chunks,
        )
        await workflow.memory.update_document_status(doc["id"], "processed", len(chunks))

        return {
            "message": "Document uploaded and processed",
            "document_id": doc["id"],
            "filename": file.filename,
            "chunks": len(chunks),
        }
    except Exception as e:
        await workflow.memory.update_document_status(doc["id"], "failed")
        logger.error(f"Document processing failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Processing failed: {e}")


def _extract_text(content: bytes, ext: str) -> str:
    """Extract plain text from uploaded file bytes."""
    if ext in {".txt", ".md"}:
        return content.decode("utf-8", errors="replace")

    # For PDF / DOCX / PPTX — attempt lightweight extraction
    if ext == ".pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            return "\n".join(page.get_text() for page in doc)
        except ImportError:
            # Fallback: just decode raw bytes (lossy but non-blocking)
            return content.decode("utf-8", errors="replace")

    if ext == ".docx":
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            return content.decode("utf-8", errors="replace")

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        except ImportError:
            return content.decode("utf-8", errors="replace")

    return content.decode("utf-8", errors="replace")


def _chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
    """Split text into overlapping chunks."""
    if not text.strip():
        return []
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks


@app.get("/documents", tags=["Documents"])
async def list_documents(user: dict = Depends(get_current_user)):
    """List all documents for the current user."""
    docs = await workflow.memory.get_documents(user["user_id"])
    return {"documents": docs}


@app.get("/documents/{document_id}", tags=["Documents"])
async def get_document(document_id: str, user: dict = Depends(get_current_user)):
    """Get a single document's metadata."""
    doc = await workflow.memory.get_document(document_id)
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    return {"document": doc}


@app.delete("/documents/{document_id}", tags=["Documents"])
async def delete_document(document_id: str, user: dict = Depends(get_current_user)):
    """Delete a document and its associated chunks."""
    doc = await workflow.memory.get_document(document_id)
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    # Mark as deleted in DB
    await workflow.memory.update_document_status(document_id, "deleted")
    return {"message": "Document deleted", "document_id": document_id}


# ═══════════════════════════════════════════════════════════════════════════════
# Routes — Tasks & Reminders
# ═══════════════════════════════════════════════════════════════════════════════

class TaskCreateRequest(BaseModel):
    title: str = Field(..., min_length=1)
    description: str = ""
    priority: str = "medium"
    due_date: str = ""


@app.post("/tasks", tags=["Tasks"], status_code=201)
async def create_task(body: TaskCreateRequest, user: dict = Depends(get_current_user)):
    """Create a new task."""
    task = await workflow.memory.create_task(
        user_id=user["user_id"],
        title=body.title,
        description=body.description,
        priority=body.priority,
        due_date=body.due_date if body.due_date else None,
    )
    return {"task": task}


@app.get("/tasks", tags=["Tasks"])
async def list_tasks(
    task_status: str = "",
    user: dict = Depends(get_current_user),
):
    """List tasks, optionally filtered by status."""
    tasks = await workflow.memory.get_tasks(user["user_id"], status=task_status)
    return {"tasks": tasks}


class ReminderCreateRequest(BaseModel):
    title: str = Field(..., min_length=1)
    message: str = ""
    remind_at: str  # ISO datetime


@app.post("/reminders", tags=["Reminders"], status_code=201)
async def create_reminder(body: ReminderCreateRequest, user: dict = Depends(get_current_user)):
    """Create a reminder."""
    reminder = await workflow.memory.create_reminder(
        user_id=user["user_id"],
        title=body.title,
        message=body.message,
        remind_at=body.remind_at,
    )
    return {"reminder": reminder}


@app.get("/reminders", tags=["Reminders"])
async def list_reminders(user: dict = Depends(get_current_user)):
    """List all reminders for the current user."""
    reminders = await workflow.memory.get_reminders(user["user_id"])
    return {"reminders": reminders}


# ═══════════════════════════════════════════════════════════════════════════════
# Routes — Habits
# ═══════════════════════════════════════════════════════════════════════════════

class HabitCreateRequest(BaseModel):
    name: str = Field(..., min_length=1)
    frequency: str = "daily"
    target_count: int = 1


@app.post("/habits", tags=["Habits"], status_code=201)
async def create_habit(body: HabitCreateRequest, user: dict = Depends(get_current_user)):
    """Create a new habit to track."""
    habit = await workflow.memory.create_habit(
        user_id=user["user_id"],
        name=body.name,
        frequency=body.frequency,
        target_count=body.target_count,
    )
    return {"habit": habit}


@app.get("/habits", tags=["Habits"])
async def list_habits(user: dict = Depends(get_current_user)):
    """List all habits for the current user."""
    habits = await workflow.memory.get_habits(user["user_id"])
    return {"habits": habits}


class HabitLogRequest(BaseModel):
    notes: str = ""


@app.post("/habits/{habit_id}/log", tags=["Habits"])
async def log_habit(habit_id: str, body: HabitLogRequest, user: dict = Depends(get_current_user)):
    """Log a habit completion."""
    result = await workflow.memory.log_habit(habit_id, body.notes)
    return {"log": result}


# ═══════════════════════════════════════════════════════════════════════════════
# Routes — Flashcards, Quizzes, Study Plans
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/flashcards", tags=["Learning"])
async def list_flashcard_sets(user: dict = Depends(get_current_user)):
    """List all flashcard sets."""
    sets = await workflow.memory.get_flashcard_sets(user["user_id"])
    return {"flashcard_sets": sets}


@app.get("/flashcards/{set_id}", tags=["Learning"])
async def get_flashcards(set_id: str, user: dict = Depends(get_current_user)):
    """Get all flashcards in a set."""
    cards = await workflow.memory.get_flashcards(set_id)
    return {"flashcards": cards}


@app.get("/quizzes", tags=["Learning"])
async def list_quizzes(user: dict = Depends(get_current_user)):
    """List all quizzes."""
    quizzes = await workflow.memory.get_quizzes(user["user_id"])
    return {"quizzes": quizzes}


@app.get("/study-plans", tags=["Learning"])
async def list_study_plans(user: dict = Depends(get_current_user)):
    """List all study plans."""
    plans = await workflow.memory.get_study_plans(user["user_id"])
    return {"study_plans": plans}


# ═══════════════════════════════════════════════════════════════════════════════
# Routes — Conversation History
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/conversations", tags=["Memory"])
async def get_conversations(
    limit: int = 20,
    user: dict = Depends(get_current_user),
):
    """Retrieve recent conversation history."""
    history = await workflow.memory.get_conversation_history(user["user_id"], limit=limit)
    return {"conversations": history}


# ═══════════════════════════════════════════════════════════════════════════════
# Global exception handler
# ═══════════════════════════════════════════════════════════════════════════════

@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    logger.error(f"Unhandled exception: {exc}", exc_info=True)
    return JSONResponse(
        status_code=500,
        content={"status": "error", "detail": "An internal server error occurred"},
    )
